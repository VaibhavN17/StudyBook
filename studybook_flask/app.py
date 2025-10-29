from flask import Flask, render_template, request, send_file
import os
import re
from werkzeug.utils import secure_filename
from pathlib import Path
from pptx import Presentation
from docx import Document
from openpyxl import load_workbook
import pdfplumber
from fpdf import FPDF

app = Flask(__name__)
UPLOAD_FOLDER = "uploads"
OUTPUT_FOLDER = "output"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

# ------------ Helper functions ------------

def extract_text_from_pdf(path):
    text = ""
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text

def extract_text_from_docx(path):
    doc = Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_pptx(path):
    prs = Presentation(path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def extract_text_from_xlsx(path):
    wb = load_workbook(path, data_only=True)
    text = ""
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            for cell in row:
                if cell:
                    text += str(cell) + " "
            text += "\n"
    return text

def extract_text(file_path):
    ext = file_path.suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    elif ext == ".xlsx":
        return extract_text_from_xlsx(file_path)
    elif ext == ".txt":
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    else:
        return ""

def clean_text(text):
    text = re.sub(r"\s+", " ", text)
    text = re.sub(r"[^A-Za-z0-9.,;:!?()\-'\s]", "", text)
    return text.strip()

def create_pdf(text, output_path):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font("Arial", size=12)
    lines = text.split(". ")
    for line in lines:
        pdf.multi_cell(0, 10, line.strip() + ".")
        pdf.ln()
    pdf.output(output_path)

# ------------ Routes ------------

@app.route("/", methods=["GET", "POST"])
def upload_files():
    if request.method == "POST":
        files = request.files.getlist("files")
        combined_text = ""
        for f in files:
            if f.filename:
                filename = secure_filename(f.filename)
                path = Path(UPLOAD_FOLDER) / filename
                f.save(path)
                text = extract_text(path)
                combined_text += "\n" + clean_text(text)
        if not combined_text.strip():
            return "No readable content found in uploaded files."
        output_path = Path(OUTPUT_FOLDER) / "StudyBook.pdf"
        create_pdf(combined_text, output_path)
        return send_file(output_path, as_attachment=True)
    return render_template("upload.html")

if __name__ == "__main__":
    app.run(debug=True)
